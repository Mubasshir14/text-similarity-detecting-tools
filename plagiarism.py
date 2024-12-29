from flask import Flask, render_template, request
import os
import re
import PyPDF2
from docx import Document
from pptx import Presentation

app = Flask(__name__)

# Set the upload folder
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Define the allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx', 'txt'}

# Create the upload folder if it doesn't exist
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def read_text_from_docx(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text.strip()

def read_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfFileReader(f)
        num_pages = reader.numPages
        for page_num in range(num_pages):
            page = reader.getPage(page_num)
            text += page.extractText()
    return text.strip()


def read_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def read_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()
    return text.strip()

def read_text_from_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    if file_extension.lower() == '.pdf':
        return read_text_from_pdf(file_path)
    elif file_extension.lower() == '.docx':
        return read_text_from_docx(file_path)
    elif file_extension.lower() == '.pptx':
        return read_text_from_pptx(file_path)
    elif file_extension.lower() == '.txt':
        return read_text_from_txt(file_path)
    else:
        return ""

def tokenize_sentences(text):
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?|\!)\s', text)
    sentences = [sentence.strip() for sentence in sentences if len(sentence.split()) > 1]
    return sentences

def calculate_similarity(text1, text2):
    sentences1 = tokenize_sentences(text1)
    sentences2 = tokenize_sentences(text2)

    similar_lines = {}  # Using a dictionary to store similar sentences and their occurrences
    total_occurrences = 0  # Total count of occurrences across all target files
    for sentence1 in sentences1:
        for sentence2 in sentences2:
            if sentence1.strip() == sentence2.strip():
                if sentence1.strip() in similar_lines:
                    similar_lines[sentence1.strip()] += 1
                else:
                    similar_lines[sentence1.strip()] = 1
                total_occurrences += 1

    total_sentences = len(sentences2)
    matching_sentences = len(similar_lines)

    similarity_percentage = (matching_sentences / total_sentences) * 100 if total_sentences > 0 else 0

    return similarity_percentage, similar_lines, total_occurrences

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/compare', methods=['POST'])
def compare():
    if 'source_file' not in request.files:
        return "No source file uploaded"

    source_file = request.files['source_file']

    if source_file.filename == '':
        return "No selected source file"

    if source_file and allowed_file(source_file.filename):
        source_file_path = os.path.join(app.config['UPLOAD_FOLDER'], source_file.filename)
        source_file.save(source_file_path)
        source_file = read_text_from_file(source_file_path)

        num_target_files = 5
        similarity_results = []
        total_similarity = 0  # Initialize total similarity

        for i in range(num_target_files):
            target_file = request.files.get(f'target_file_{i}')
            if target_file and allowed_file(target_file.filename):
                target_file_path = os.path.join(app.config['UPLOAD_FOLDER'], target_file.filename)
                target_file.save(target_file_path)
                target_file = read_text_from_file(target_file_path)
                similarity_percentage, similar_lines, file_occurrences = calculate_similarity(target_file, source_file)
                similarity_results.append((f"Target File {i + 1}", similarity_percentage, similar_lines, file_occurrences))
                total_similarity += similarity_percentage  # Add similarity percentage to total
            else:
                similarity_results.append((f"Target File {i + 1}", "Invalid file format", {}, 0))

        average_similarity = total_similarity / num_target_files  # Calculate average similarity

        print("Average Similarity:", average_similarity)  # Print average similarity

        return render_template('results.html', similarity_results=similarity_results, average_similarity=average_similarity)
    else:
        return "Invalid file format for source file"

if __name__ == "__main__":
    app.run(debug=True)
